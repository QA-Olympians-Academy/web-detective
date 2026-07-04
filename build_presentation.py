"""
Builds web-detective-workshop.pptx  (~100 slides, 8-hour workshop deck).
Run:  python3 build_presentation.py
"""
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Dimensions ────────────────────────────────────────────────────────────────
W  = 12188952   # slide width  (EMU)
H  = 6858000    # slide height (EMU)

# ── Colour palette ────────────────────────────────────────────────────────────
def c(h): return RGBColor(int(h[:2],16), int(h[2:4],16), int(h[4:],16))
BG      = c('0D1117')
CARD    = c('161B22')
CARD2   = c('1C2128')
WHITE   = c('FFFFFF')
MUTED   = c('8B949E')
BLACK   = c('0D1117')
BLUE    = c('58A6FF')
PURPLE  = c('D2A8FF')
GREEN   = c('3FB950')
BGREEN  = c('39D353')
ORANGE  = c('F78466')
HORANGE = c('FFA657')
YELLOW  = c('E3B341')
RED     = c('FF7B72')

CH_COLOR = {0:BLUE, 1:PURPLE, 2:BLUE, 3:GREEN, 4:ORANGE, 4.1:HORANGE,
            5:BGREEN, 6:PURPLE, 7:BLUE}

# ── Low-level helpers ─────────────────────────────────────────────────────────

def rect(slide, l, t, w, h, fill=None):
    s = slide.shapes.add_shape(1, Emu(l), Emu(t), Emu(w), Emu(h))
    s.line.fill.background()
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else:    s.fill.background()
    return s

def tb(slide, l, t, w, h, text, sz, color=WHITE, bold=False,
       align=PP_ALIGN.LEFT, wrap=True, font=None):
    box = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    box.text_frame.word_wrap = wrap
    p = box.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(sz)
    run.font.bold = bold
    run.font.color.rgb = color
    if font: run.font.name = font
    return box

def multiline(slide, l, t, w, h, lines, sz, default_color=WHITE,
              bold=False, align=PP_ALIGN.LEFT, font=None, line_spacing=None):
    """lines = list of str or (str, color)"""
    box = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    box.text_frame.word_wrap = True
    tf = box.text_frame
    for i, line in enumerate(lines):
        text, clr = (line, default_color) if isinstance(line, str) else line
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        if line_spacing:
            from pptx.util import Pt as _Pt
            from pptx.oxml.ns import qn
            from lxml import etree
            pPr = para._p.get_or_add_pPr()
            lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
            spcPts = etree.SubElement(lnSpc, qn('a:spcPts'))
            spcPts.set('val', str(int(line_spacing * 100)))
        run = para.add_run()
        run.text = text
        run.font.size = Pt(sz)
        run.font.bold = bold
        run.font.color.rgb = clr
        if font: run.font.name = font
    return box

def bg(slide): rect(slide, 0, 0, W, H, BG)

def topbar(slide, color): rect(slide, 0, 0, W, 50292, color)

def chapbadge(slide, label, color):
    rect(slide, 274320, 137160, 1371600, 347472, color)
    tb(slide, 329184, 164592, 1280160, 320040, label, 13,
       BLACK, bold=True, align=PP_ALIGN.CENTER)

def chapter_label_top(slide, label, color):
    tb(slide, 365760, 80000, 3000000, 280000, label, 11, color)

# ── Slide type builders ───────────────────────────────────────────────────────

def slide_title(slide, title, sub1, sub2, tech_line):
    """Opening title slide."""
    bg(slide)
    rect(slide, 0, 0, W, 50292, BLUE)
    rect(slide, 0, H - 50292, W, 50292, BLUE)
    # Decorative vertical accent
    rect(slide, 0, 50292, 12000, H - 100584, BLUE)
    tb(slide, 200000, 1400000, W - 400000, 700000,
       title, 42, WHITE, bold=True, align=PP_ALIGN.LEFT)
    tb(slide, 200000, 2180000, W - 400000, 500000,
       sub1, 22, BLUE, bold=True)
    tb(slide, 200000, 2740000, W - 400000, 400000,
       sub2, 15, MUTED)
    tb(slide, 200000, 3600000, W - 400000, 300000,
       tech_line, 13, MUTED)

def slide_chapter_opener(slide, ch_label, time_str, title, subtitle, color):
    bg(slide)
    topbar(slide, color)
    chapbadge(slide, f"{ch_label}  ·  {time_str}", color)
    tb(slide, 1874519, 109728, 9900000, 600000, title, 30, WHITE, bold=True)
    tb(slide, 1874519, 660000, 9900000, 347472, subtitle, 15, MUTED)

def slide_concept(slide, ch_label, color, title, bullets,
                  col2_title=None, col2_bullets=None):
    """Single or two-column concept slide."""
    bg(slide)
    topbar(slide, color)
    chapter_label_top(slide, ch_label, color)
    tb(slide, 365760, 140000, W - 730000, 500000, title, 24, WHITE, bold=True)

    if col2_title is None:
        # single card
        card_top = 680000; card_h = H - card_top - 120000
        rect(slide, 274320, card_top, W - 548640, card_h, CARD)
        multiline(slide, 502920, card_top + 120000, W - 1005840,
                  card_h - 200000, bullets, 14, line_spacing=18)
    else:
        # two cards
        card_top = 680000; card_h = H - card_top - 120000
        half = (W - 548640 - 91440) // 2
        rect(slide, 274320,             card_top, half, card_h, CARD)
        rect(slide, 274320 + half + 91440, card_top, half, card_h, CARD)
        multiline(slide, 502920, card_top + 120000,
                  half - 228600, card_h - 200000, bullets, 13, line_spacing=17)
        tb(slide, 274320 + half + 91440 + 182880, card_top + 120000,
           half - 228600, 300000, col2_title, 14, color, bold=True)
        multiline(slide, 274320 + half + 91440 + 182880,
                  card_top + 480000, half - 228600,
                  card_h - 580000, col2_bullets, 13, line_spacing=17)

def slide_code(slide, ch_label, color, title, filename, code_lines):
    """Code slide: header + dark card with monospace code."""
    bg(slide)
    topbar(slide, color)
    chapter_label_top(slide, ch_label, color)
    tb(slide, 365760, 140000, 7000000, 440000, title, 22, WHITE, bold=True)
    tb(slide, 7500000, 200000, 4500000, 320000, filename, 12, MUTED,
       font="Courier New")
    card_top = 640000; card_h = H - card_top - 120000
    rect(slide, 274320, card_top, W - 548640, card_h, CARD)
    multiline(slide, 480000, card_top + 100000, W - 960000,
              card_h - 180000, code_lines, 12, BGREEN,
              font="Courier New", line_spacing=15)

def slide_two_panel(slide, ch_label, color, title,
                    left_title, left_items,
                    right_title, right_items):
    bg(slide)
    topbar(slide, color)
    chapter_label_top(slide, ch_label, color)
    tb(slide, 365760, 140000, W - 730000, 440000, title, 22, WHITE, bold=True)
    card_top = 640000; card_h = H - card_top - 120000
    half = (W - 548640 - 91440) // 2
    Lx = 274320; Rx = Lx + half + 91440
    rect(slide, Lx, card_top, half, card_h, CARD)
    rect(slide, Rx, card_top, half, card_h, CARD)
    tb(slide, Lx+182880, card_top+100000, half-228600, 300000,
       left_title, 14, color, bold=True)
    multiline(slide, Lx+182880, card_top+450000, half-228600,
              card_h-550000, left_items, 13, line_spacing=17)
    tb(slide, Rx+182880, card_top+100000, half-228600, 300000,
       right_title, 14, color, bold=True)
    multiline(slide, Rx+182880, card_top+450000, half-228600,
              card_h-550000, right_items, 13, line_spacing=17)

def slide_table(slide, ch_label, color, title, headers, rows):
    bg(slide)
    topbar(slide, color)
    chapter_label_top(slide, ch_label, color)
    tb(slide, 365760, 140000, W - 730000, 440000, title, 22, WHITE, bold=True)
    n_cols = len(headers)
    col_w = (W - 548640) // n_cols
    row_h = 420000
    table_top = 680000
    # header row
    for ci, h in enumerate(headers):
        rect(slide, 274320 + ci*col_w, table_top, col_w - 18000, row_h, color)
        tb(slide, 274320 + ci*col_w + 60000, table_top + 120000,
           col_w - 120000, row_h - 60000, h, 13, BLACK, bold=True)
    for ri, row in enumerate(rows):
        ry = table_top + (ri+1)*row_h + ri*18000
        fill = CARD if ri % 2 == 0 else CARD2
        for ci, cell in enumerate(row):
            rect(slide, 274320 + ci*col_w, ry, col_w - 18000, row_h, fill)
            clr = row[ci][1] if isinstance(row[ci], tuple) else WHITE
            txt = row[ci][0] if isinstance(row[ci], tuple) else row[ci]
            tb(slide, 274320 + ci*col_w + 60000, ry + 120000,
               col_w - 120000, row_h - 60000, txt, 12, clr,
               font="Courier New" if ci == 0 else None)

def slide_discussion(slide, ch_label, color, questions):
    bg(slide)
    topbar(slide, color)
    chapter_label_top(slide, ch_label, color)
    tb(slide, 365760, 140000, W - 730000, 440000,
       "Discussion Questions", 24, WHITE, bold=True)
    q_top = 700000; q_gap = (H - q_top - 200000) // 3
    for i, q in enumerate(questions):
        rect(slide, 274320, q_top + i*q_gap, W - 548640, q_gap - 60000, CARD)
        tb(slide, 365760, q_top + i*q_gap + 60000, 700000, q_gap - 120000,
           str(i+1), 28, color, bold=True, align=PP_ALIGN.CENTER)
        tb(slide, 1100000, q_top + i*q_gap + 80000, W - 1400000,
           q_gap - 160000, q, 14, WHITE)

def slide_tasks(slide, ch_label, color, tasks):
    bg(slide)
    topbar(slide, color)
    chapter_label_top(slide, ch_label, color)
    tb(slide, 365760, 140000, W - 730000, 440000,
       "Workshop Tasks", 24, WHITE, bold=True)
    t_top = 680000; t_h = (H - t_top - 180000) // len(tasks)
    for i, (task_title, task_body) in enumerate(tasks):
        ty = t_top + i * t_h
        rect(slide, 274320, ty, W - 548640, t_h - 60000, CARD)
        tb(slide, 365760, ty + 60000, 800000, t_h - 120000,
           f"Task {chr(65+i)}", 13, color, bold=True)
        tb(slide, 1200000, ty + 60000, 3600000, t_h - 120000,
           task_title, 14, WHITE, bold=True)
        tb(slide, 4900000, ty + 60000, W - 5200000, t_h - 120000,
           task_body, 12, MUTED)

def slide_break(slide, label, time_range, color=BLUE):
    bg(slide)
    rect(slide, 0, 0, W, H//3, color)
    rect(slide, 0, H//3, W, H - H//3, CARD)
    tb(slide, 0, H//3 - 200000, W, 400000, label, 36, BLACK,
       bold=True, align=PP_ALIGN.CENTER)
    tb(slide, 0, H//3 + 280000, W, 400000, time_range, 20, MUTED,
       align=PP_ALIGN.CENTER)

def slide_full_text(slide, big_text, sub_text, color=BLUE):
    bg(slide)
    topbar(slide, color)
    rect(slide, 274320, 50292, W - 548640, H - 50292, CARD)
    tb(slide, 548640, 1800000, W - 1097280, 1200000, big_text,
       32, WHITE, bold=True, align=PP_ALIGN.CENTER)
    tb(slide, 548640, 3200000, W - 1097280, 600000, sub_text,
       16, MUTED, align=PP_ALIGN.CENTER)

def slide_takeaways(slide, items):
    bg(slide)
    topbar(slide, BLUE)
    tb(slide, 365760, 100000, W - 730000, 500000,
       "Key Takeaways", 30, WHITE, bold=True)
    row_h = 680000
    start = 640000
    for i, (headline, detail) in enumerate(items[:7]):
        rx = 274320 + (i % 2) * (W // 2)
        ry = start + (i // 2) * row_h
        if i == 6:  # last one centered
            rx = (W - W // 2) // 2 + 100000
        rw = W // 2 - 320000
        rect(slide, rx, ry, rw, row_h - 60000, CARD)
        tb(slide, rx + 120000, ry + 60000, rw - 240000, 240000,
           headline, 14, BLUE, bold=True)
        tb(slide, rx + 120000, ry + 320000, rw - 240000, 300000,
           detail, 12, MUTED)

def add_slide(prs, color=None):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    for shape in list(s.shapes):
        shape.element.getparent().remove(shape.element)
    return s

# ── Content data ──────────────────────────────────────────────────────────────

SLIDE_W = W

def build(prs):
    s = lambda: add_slide(prs)

    # ── 1. Title ──────────────────────────────────────────────────────────────
    sl = s()
    slide_title(sl,
        "Web Detective",
        "AI-Driven Test Automation Workshop",
        "Building autonomous agents that test modern web apps\n"
        "with LLMs, browser automation & structured reasoning",
        "Playwright  ·  Ollama + DeepSeek-R1  ·  MCP  ·  GitHub Actions  ·  8 hours  ·  Senior QA / SDET")

    # ── 2. What You Will Build ────────────────────────────────────────────────
    sl = s()
    slide_concept(sl, "Welcome", BLUE, "What You Will Build Today",
        [("By the end of this workshop you will have shipped:", MUTED),
         ("✓  A self-healing test suite — selectors that auto-repair at runtime", WHITE),
         ("✓  A custom AI agent that observes, plans, acts, and verifies", WHITE),
         ("✓  An MCP server exposing your browser as structured tools", WHITE),
         ("✓  AI-generated test plans converted to runnable Playwright specs", WHITE),
         ("✓  A CI/CD pipeline that annotates failures inline in PR diffs", WHITE),
         ("✓  Reusable testing skills invokable from Claude Code or VS Code", WHITE),
         ("", WHITE),
         ("Everything runs against a real React app — no mocks, no shortcuts.", MUTED)])

    # ── 3. Agenda ─────────────────────────────────────────────────────────────
    sl = s()
    bg(sl); topbar(sl, BLUE)
    tb(sl, 365760, 100000, W - 730000, 500000,
       "Workshop Agenda", 28, WHITE, bold=True)
    tb(sl, 365760, 620000, W - 730000, 300000,
       "8-hour deep dive into production-grade agentic test automation", 14, MUTED)
    rows_data = [
        ("Ch1",   "09:30", "Foundations of Agentic Automation",              "75 min",  PURPLE),
        ("Ch2",   "11:00", "Browser as Execution Layer + CLI Deep-Dive",      "90 min",  BLUE),
        ("Ch3",   "13:15", "Build Your Own MCP Server",                       "60 min",  GREEN),
        ("Ch4",   "14:15", "Self-Healing Selectors",                          "45 min",  ORANGE),
        ("Ch4.1", "15:00", "Playwright Agents: Planner, Generator, Healer",   "30 min",  HORANGE),
        ("Ch5",   "15:45", "Custom AI Agent (Reasoning Loop)",                "45 min",  BGREEN),
        ("Ch6",   "16:30", "Creating Effective Testing Skills",               "30 min",  PURPLE),
        ("Ch7",   "17:00", "Agent & CI/CD Pipeline",                         "60 min",  BLUE),
    ]
    row_h = 640000; start_y = 980000
    for i, (lbl, t, title, dur, col) in enumerate(rows_data):
        ry = start_y + i * row_h
        rect(sl, 274320,          ry, W - 548640, row_h - 50000, CARD)
        rect(sl, 274320,          ry, 60000,       row_h - 50000, col)
        tb(sl, 440000,   ry + 90000, 700000,  row_h - 180000, lbl,   13, col, bold=True)
        tb(sl, 1150000,  ry + 90000, 900000,  row_h - 180000, t,     12, MUTED)
        tb(sl, 2150000,  ry + 90000, 7500000, row_h - 180000, title, 15, WHITE, bold=True)
        tb(sl, 10600000, ry + 90000, 1300000, row_h - 180000, dur,   12, MUTED, align=PP_ALIGN.RIGHT)

    # ── 4. Prerequisites ──────────────────────────────────────────────────────
    sl = s()
    slide_table(sl, "Setup", BLUE, "Prerequisites — Install Before Arriving",
        ["Tool", "Purpose", "Verify"],
        [["VS Code + GitHub Copilot",    "IDE + Agent Mode",              "code --version"],
         ["Claude Code (AI CLI)",         "Skills + agentic authoring",    "claude --version"],
         ["Node.js LTS",                  "Runtime for Playwright / TS",   "node -v"],
         ["Playwright",                   "Browser automation + MCP",      "npx playwright --version"],
         ["Docker Engine",                "CI/CD containerised runs",      "docker info"],
         ["Ollama + DeepSeek-R1",         "Local LLM (no API key)",        "ollama pull deepseek-r1:8b"]])

    # ── 5. System Under Test ──────────────────────────────────────────────────
    sl = s()
    slide_two_panel(sl, "Setup", BLUE, "The System Under Test — Web Detective",
        "Tech Stack",
        [("React 18 + TypeScript + Vite", WHITE),
         ("React Router v6", WHITE),
         ("Recharts (dashboard charts)", WHITE),
         ("Playwright (automation layer)", WHITE),
         ("Ollama + DeepSeek-R1 (local LLM reasoning)", WHITE),
         ("", WHITE),
         ("npm install  &&  npm run dev", BGREEN),
         ("→  http://localhost:5173", MUTED)],
        "App Routes",
        [("/login   Auth form, session redirect", WHITE),
         ("/dashboard   Stats + Recharts charts", WHITE),
         ("/products   Searchable product table", WHITE),
         ("", WHITE),
         ("Credentials", MUTED),
         ("admin@shop.com  /  password123", BGREEN),
         ("", WHITE),
         ("Purpose: intentionally simple so focus\nstays on automation patterns.", MUTED)])

    # ── 6. Repo Structure ─────────────────────────────────────────────────────
    sl = s()
    slide_code(sl, "Setup", BLUE, "Repository Structure", "tree -L 2",
        [("src/                  React app (the system under test)", MUTED),
         ("tests/                Hand-written Playwright specs + POMs", MUTED),
         ("  fixtures/           Shared setup (auth, baseURL)", MUTED),
         ("  pages/              Page Object Models", MUTED),
         ("examples/             Workshop chapter code", MUTED),
         ("  ch1-foundations/    brittle-test.spec.ts, task-graph.ts", WHITE),
         ("  ch2-execution-layer/ action-wrapper.ts, browser-session.ts", WHITE),
         ("  ch3-mcp/            server.ts, playwright-mcp-client.ts", WHITE),
         ("  ch4-self-healing/   locator-store.ts, self-healer.ts", WHITE),
         ("  ch4.1-playwright-agents/  planner.ts, healer.ts", HORANGE),
         ("  ch5-custom-agent/   agent.ts, tools.ts, prompts.ts", WHITE),
         ("  ch7-agent-ci/       agent-runner.ts, reporter.ts, scenarios.ts", WHITE),
         (".claude/skills/       pw-run, pw-debug, pw-plan, pw-self-heal …", PURPLE),
         (".github/workflows/    per-chapter CI jobs", MUTED)])

    # ── 7. Environment Check ──────────────────────────────────────────────────
    sl = s()
    slide_concept(sl, "Setup", BLUE, "Environment Check — Run These Now",
        [("npx playwright test --reporter=list", BGREEN),
         ("  → should show tests/auth.spec.ts, products.spec.ts, dashboard.spec.ts", MUTED),
         ("", WHITE),
         ("ollama run deepseek-r1:8b 'hello'", BGREEN),
         ("  → a response = local model wired correctly", MUTED),
         ("", WHITE),
         ("npx ts-node examples/ch5-custom-agent/agent.ts", BGREEN),
         ("  → agent should log steps and pass the full e-commerce scenario", MUTED),
         ("", WHITE),
         ("If anything fails → raise your hand now.", ORANGE)])

    # ── 8. Ch1 Opener ─────────────────────────────────────────────────────────
    sl = s()
    slide_chapter_opener(sl, "CH 1", "09:30",
        "Foundations of Agentic Automation",
        "75 min  ·  Why traditional scripts fail and how agents fix it", PURPLE)

    # ── 9. What "Agentic" Means ───────────────────────────────────────────────
    sl = s()
    slide_two_panel(sl, "CH 1", PURPLE, "Traditional Script vs. Agentic System",
        "Traditional Playwright Test",
        [("Hardcoded sequence of steps", WHITE),
         ("Breaks on any selector change", RED),
         ("No awareness of page state", RED),
         ("No recovery from failure", RED),
         ("Knows exactly one path to success", RED),
         ("Requires human to update after\nevery UI change", RED)],
        "Agentic System",
        [("Goal-driven: given a task, figures out steps", WHITE),
         ("Observes current state before acting", BGREEN),
         ("Adapts: tries alternatives on failure", BGREEN),
         ("Verifies: checks success criteria", BGREEN),
         ("Learns: updates locator memory", BGREEN),
         ("Self-heals: repairs broken selectors\nwithout human intervention", BGREEN)])

    # ── 10. The 5-Phase Agentic Loop ──────────────────────────────────────────
    sl = s()
    bg(sl); topbar(sl, PURPLE); chapter_label_top(sl, "CH 1", PURPLE)
    tb(sl, 365760, 140000, W - 730000, 440000,
       "The Agentic Reasoning Loop", 24, WHITE, bold=True)
    phases = [
        ("① Observe",  "ariaSnapshot() — capture page state as structured text"),
        ("② Plan",     "Decompose goal into steps; choose next action"),
        ("③ Act",      "click() / fill() / navigate() — modify browser state"),
        ("④ Verify",   "assert_visible() / assert_text() — confirm expectations"),
        ("⑤ Learn",    "Store healed selectors; log timing; update memory"),
    ]
    ph_w = (W - 548640) // 5 - 30000
    ph_top = 680000
    for i, (phase, desc) in enumerate(phases):
        px = 274320 + i * ((W - 548640) // 5 + 6000)
        rect(sl, px, ph_top, ph_w, H - ph_top - 200000, CARD)
        rect(sl, px, ph_top, ph_w, 80000, PURPLE)
        tb(sl, px + 40000, ph_top + 120000, ph_w - 80000, 400000,
           phase, 15, PURPLE, bold=True, align=PP_ALIGN.CENTER)
        tb(sl, px + 40000, ph_top + 560000, ph_w - 80000,
           H - ph_top - 900000, desc, 12, MUTED,
           align=PP_ALIGN.CENTER)
        if i < 4:
            tb(sl, px + ph_w + 6000, ph_top + (H - ph_top - 200000)//2 - 100000,
               30000, 200000, "→", 18, PURPLE, bold=True)

    # ── 11. Architectural Patterns ────────────────────────────────────────────
    sl = s()
    bg(sl); topbar(sl, PURPLE); chapter_label_top(sl, "CH 1", PURPLE)
    tb(sl, 365760, 140000, W - 730000, 440000,
       "Three Agentic Architectural Patterns", 24, WHITE, bold=True)
    patterns = [
        ("Reactive Agent", PURPLE,
         "Responds to specific triggers.\nEvent fires → observe → act → done.\nExample: healer fires when a locator fails.",
         "Simple, predictable,\nlimited to predefined triggers."),
        ("Goal-Driven Agent", BLUE,
         "Given a high-level goal, plans and executes a sequence of actions.\nExample: WebTestAgent.run(tasks.fullEcommerce())",
         "Flexible, adapts to page state,\nmay take multiple paths to succeed."),
        ("Multi-Agent Pipeline", BGREEN,
         "Specialist agents hand off to each other:\nPlanner → Generator → Healer.\nExample: Ch4.1 Playwright Agents.",
         "Each agent is small and focused;\npipeline handles complex workflows."),
    ]
    col_w = (W - 548640 - 180000) // 3
    col_top = 680000
    for i, (title, col, desc, tradeoff) in enumerate(patterns):
        cx = 274320 + i * (col_w + 90000)
        rect(sl, cx, col_top, col_w, H - col_top - 200000, CARD)
        rect(sl, cx, col_top, col_w, 60000, col)
        tb(sl, cx + 120000, col_top + 120000, col_w - 240000, 300000,
           title, 16, col, bold=True)
        tb(sl, cx + 120000, col_top + 500000, col_w - 240000,
           H - col_top - 1100000, desc, 13, WHITE)
        rect(sl, cx + 60000, H - 700000, col_w - 120000, 380000, CARD2)
        tb(sl, cx + 120000, H - 640000, col_w - 240000, 320000,
           tradeoff, 12, MUTED)

    # ── 12. Brittle Test Anti-Patterns ────────────────────────────────────────
    sl = s()
    slide_code(sl, "CH 1", PURPLE, "Anti-Patterns in brittle-test.spec.ts",
               "ch1-foundations/brittle-test.spec.ts",
        [("await page.waitForTimeout(2000)         // ❌ masks timing, not fixes it", RED),
         ("", WHITE),
         ("await page.click(", WHITE),
         ("  '.login-card > form > div:nth-child(1) > input'  // ❌ structural CSS", RED),
         (")", WHITE),
         ("", WHITE),
         ("await page.click('.login-card > form > button:nth-child(1)')  // ❌ position", RED),
         ("", WHITE),
         ("await page.waitForTimeout(1500)         // ❌ another sleep", RED),
         ("", WHITE),
         ("const heading = await page.$('.page-header > h2')  // ❌ CSS class assert", RED),
         ("const text = await heading?.textContent()", WHITE),
         ("expect(text).toBe('Dashboard')", WHITE),
         ("", WHITE),
         ("// Resilient alternative:", MUTED),
         ("await page.getByRole('heading', { name: 'Dashboard' }).waitFor()  // ✓", BGREEN)])

    # ── 13. Resilient vs Brittle ──────────────────────────────────────────────
    sl = s()
    slide_two_panel(sl, "CH 1", PURPLE, "From Brittle to Resilient",
        "Brittle (avoid)",
        [("page.waitForTimeout(2000)", RED),
         ("→ Replace with: page.waitForLoadState()\n   or element.waitFor({ state: 'visible' })", MUTED),
         ("", WHITE),
         ("'.login-card > form > div:nth-child(1)'", RED),
         ("→ Replace with:\n   page.getByLabel('Email address')", MUTED),
         ("", WHITE),
         ("$eval('.form-error', el => el.textContent)", RED),
         ("→ Replace with:\n   expect(page.getByText('Invalid')).toBeVisible()", MUTED)],
        "Resilient (prefer)",
        [("ARIA-first selector hierarchy:", WHITE),
         ("1. getByRole('button', { name: 'Sign In' })", BGREEN),
         ("2. getByLabel('Email address')", BGREEN),
         ("3. getByPlaceholder('you@example.com')", BGREEN),
         ("4. getByText('Forgot password?')", BGREEN),
         ("5. locator('.semantic-class')", YELLOW),
         ("6. locator('[data-testid=submit]')", YELLOW),
         ("7. Structural CSS  (never use)", RED)])

    # ── 14. Task Graph Concept ────────────────────────────────────────────────
    sl = s()
    slide_concept(sl, "CH 1", PURPLE, "Task Graphs — Data-Driven Test Plans",
        [("A task graph is a DAG where each node is an AgentTask:", MUTED),
         ("", WHITE),
         ("  id:            unique key (e.g. 'submitCredentials')", WHITE),
         ("  goal:          declarative description ('verify login redirects to dashboard')", WHITE),
         ("  preconditions: tasks that must pass first", WHITE),
         ("  action:        async (page) => void", WHITE),
         ("  verify:        async (page) => boolean", WHITE),
         ("", WHITE),
         ("Why this matters:", MUTED),
         ("→ Runner checks preconditions before executing — no silent order bugs", WHITE),
         ("→ Each task has a binary pass/fail — observable, loggable, retryable", WHITE),
         ("→ Agent can skip already-passed nodes on retry — partial re-runs", WHITE),
         ("→ Graph structure is the test plan — readable by humans and LLMs", WHITE)])

    # ── 15. Task Graph Code ───────────────────────────────────────────────────
    sl = s()
    slide_code(sl, "CH 1", PURPLE, "Task Graph Implementation",
               "ch1-foundations/task-graph.ts",
        [("interface AgentTask {", WHITE),
         ("  id:            string", WHITE),
         ("  goal:          string", WHITE),
         ("  preconditions: string[]", WHITE),
         ("  action:        (page: Page) => Promise<void>", WHITE),
         ("  verify:        (page: Page) => Promise<boolean>", WHITE),
         ("}", WHITE),
         ("", WHITE),
         ("const loginTaskGraph: AgentTask[] = [", WHITE),
         ("  { id: 'navigateToLogin',", WHITE),
         ("    goal: 'Open the login page',", MUTED),
         ("    preconditions: [],", WHITE),
         ("    action:  async p => p.goto('/login'),", WHITE),
         ("    verify:  async p => p.url().includes('/login') },", WHITE),
         ("", WHITE),
         ("  { id: 'submitCredentials',", WHITE),
         ("    preconditions: ['navigateToLogin'],   // enforced by runner", HORANGE),
         ("    action: async p => { /* fill + click */ },", WHITE),
         ("    verify: async p => p.url().includes('/dashboard') },", WHITE),
         ("]", WHITE)])

    # ── 16. Ch1 Discussion ────────────────────────────────────────────────────
    sl = s()
    slide_discussion(sl, "CH 1", PURPLE, [
        "The brittle test uses await page.waitForTimeout(2000). What specific problem "
        "does this mask, and what Playwright API should replace it?",
        "Which of the five agentic phases (observe / plan / act / verify / learn) is absent "
        "from a traditional Playwright test, and why does its absence matter at scale?",
        "loginTaskGraph has three tasks with preconditions. What happens if submitCredentials "
        "runs before navigateToLogin completes — and how does the task graph prevent it?"])

    # ── 17. Ch1 Workshop Tasks ────────────────────────────────────────────────
    sl = s()
    slide_tasks(sl, "CH 1", PURPLE, [
        ("Run the brittle test",
         "npx playwright test examples/ch1-foundations/brittle-test.spec.ts\n"
         "Rename .login-card → .auth-card in Login.tsx — watch it fail for the wrong reason."),
        ("Map anti-patterns to the loop",
         "For each ❌ in brittle-test.spec.ts, identify which agentic phase it prevents\n"
         "(observe/plan/act/verify/learn). Write the ARIA-first replacement selector."),
        ("Read the resilient version",
         "Open tests/auth.spec.ts. Note every getByRole / getByLabel call.\n"
         "Count how many CSS selectors remain — aim for zero."),
        ("Extend the task graph",
         "Add a logoutTask node to loginTaskGraph with preconditions: ['submitCredentials'].\n"
         "Implement action + verify. Run the graph via the provided runner."),
    ])

    # ── 18. BREAK ─────────────────────────────────────────────────────────────
    sl = s()
    slide_break(sl, "☕  Break", "10:45 – 11:00  ·  15 minutes", BLUE)

    # ── 19. Ch2 Opener ────────────────────────────────────────────────────────
    sl = s()
    slide_chapter_opener(sl, "CH 2", "11:00",
        "Browser as Execution Layer",
        "90 min  ·  Build a fault-tolerant action toolkit that LLMs can call", BLUE)

    # ── 20. Browser as Vocabulary ─────────────────────────────────────────────
    sl = s()
    slide_concept(sl, "CH 2", BLUE, "Playwright as an Action Vocabulary — Not a Test Runner",
        [("Key mental shift:", MUTED),
         ("  Traditional mindset:  'I am writing a test.'", ORANGE),
         ("  Agentic mindset:      'I am defining tools an LLM can call.'", BGREEN),
         ("", WHITE),
         ("Every browser interaction becomes a typed tool with:", WHITE),
         ("  • A clear input schema  (selector, value, url …)", WHITE),
         ("  • A structured ActionResult  { ok, durationMs, error?, snapshot? }", WHITE),
         ("  • Retry logic built in  — no bare await page.click()", WHITE),
         ("", WHITE),
         ("This wrapper layer is what lets an LLM orchestrate Playwright", MUTED),
         ("without knowing anything about the test framework internals.", MUTED)])

    # ── 21. ActionWrapper API ─────────────────────────────────────────────────
    sl = s()
    slide_code(sl, "CH 2", BLUE, "ActionWrapper — The LLM-Callable Interface",
               "ch2-execution-layer/action-wrapper.ts",
        [("interface ActionResult {", WHITE),
         ("  ok:          boolean         // false = agent must re-plan", WHITE),
         ("  durationMs:  number          // perf signal for adaptive agents", WHITE),
         ("  error?:      string          // structured — not a raw exception", WHITE),
         ("  snapshot?:   string          // a11y tree after the action", WHITE),
         ("}", WHITE),
         ("", WHITE),
         ("class ActionWrapper {", WHITE),
         ("  async click(selector: string, desc: string): Promise<ActionResult>", BLUE),
         ("  async fill(selector: string, value: string): Promise<ActionResult>", BLUE),
         ("  async navigate(url: string):                 Promise<ActionResult>", BLUE),
         ("  async snapshot():                            Promise<ActionResult>", BLUE),
         ("  async assertVisible(selector, desc):         Promise<ActionResult>", BLUE),
         ("}", WHITE),
         ("", WHITE),
         ("// Every call returns ok:false instead of throwing.", MUTED),
         ("// The agent reads ok and decides whether to retry, re-plan, or fail.", MUTED)])

    # ── 22. BrowserSession ────────────────────────────────────────────────────
    sl = s()
    slide_code(sl, "CH 2", BLUE, "BrowserSession — Fault-Tolerant Navigation",
               "ch2-execution-layer/browser-session.ts",
        [("class BrowserSession {", WHITE),
         ("  async navigateWithRetry(url: string, maxRetries = 3) {", WHITE),
         ("    for (let attempt = 1; attempt <= maxRetries; attempt++) {", WHITE),
         ("      try {", WHITE),
         ("        await this.page.goto(url, { waitUntil: 'domcontentloaded' })", WHITE),
         ("        return  // success", BGREEN),
         ("      } catch {", WHITE),
         ("        if (attempt === maxRetries) throw", WHITE),
         ("        await new Promise(r =>", WHITE),
         ("          setTimeout(r, this.retryDelayMs * attempt))  // linear back-off", ORANGE),
         ("      }", WHITE),
         ("    }", WHITE),
         ("  }", WHITE),
         ("}", WHITE),
         ("", WHITE),
         ("// Note: retryDelayMs * attempt is LINEAR, not exponential.", MUTED),
         ("// Exponential: retryDelayMs * 2^(attempt-1)", MUTED),
         ("// Use exponential when failures signal server overload.", MUTED)])

    # ── 23. Page State Extraction ─────────────────────────────────────────────
    sl = s()
    slide_two_panel(sl, "CH 2", BLUE, "Extracting Page State for LLM Reasoning",
        "Accessibility Tree  (prefer)",
        [("page.locator('body').ariaSnapshot()", BGREEN),
         ("", WHITE),
         ("✓ Text — no image tokens", WHITE),
         ("✓ Headless — works in CI", WHITE),
         ("✓ Fast — no screenshot latency", WHITE),
         ("✓ Structured — roles, names, states", WHITE),
         ("✓ Better for forms, tables, lists", WHITE),
         ("", WHITE),
         ("Cap at 3000–4000 chars to stay\nwithin context budget.", MUTED)],
        "Screenshot  (fallback)",
        [("await page.screenshot({ path: '...' })", ORANGE),
         ("", WHITE),
         ("✓ Captures charts, canvas, SVGs", WHITE),
         ("✓ Useful for visual regression", WHITE),
         ("✗ Slow — render + encode", WHITE),
         ("✗ Expensive — image tokens", WHITE),
         ("✗ Sensitive to viewport size", WHITE),
         ("", WHITE),
         ("Use for Recharts on /dashboard\nor any canvas-based component.", MUTED)])

    # ── 24. Playwright CLI Commands ───────────────────────────────────────────
    sl = s()
    slide_table(sl, "CH 2", BLUE, "Playwright CLI — Commands Every Agent QA Must Know",
        ["Command", "Purpose"],
        [["npx playwright test --headed --debug",  "Step through test interactively in Inspector"],
         ["npx playwright codegen <url>",           "Record interactions, auto-generate ARIA selectors"],
         ["npx playwright show-report",             "HTML report: timeline, traces, screenshots"],
         ["npx playwright show-trace <file>",       "Offline trace viewer — step through every action"],
         ["PWDEBUG=1 npx playwright test",          "Launch Playwright Inspector on first test"],
         ["npx playwright test --ui",               "Interactive test runner UI with live results"]])

    # ── 25. codegen Demo ──────────────────────────────────────────────────────
    sl = s()
    slide_concept(sl, "CH 2", BLUE, "codegen — Record and Refine Locators",
        [("Run:  npx playwright codegen http://localhost:5173", BGREEN),
         ("", WHITE),
         ("Codegen records every click and fill, then suggests selectors.", MUTED),
         ("It prefers ARIA automatically — watch it choose:", WHITE),
         ("  getByLabel('Email address')    over    input[type='email']", WHITE),
         ("  getByRole('button', {name: 'Sign In'})    over    .btn-login", WHITE),
         ("", WHITE),
         ("Hands-on:", MUTED),
         ("1. Run codegen against /login — perform a full login", WHITE),
         ("2. Copy the generated test to a scratch file", WHITE),
         ("3. Count CSS selectors vs ARIA selectors — aim for 0 CSS", WHITE),
         ("4. Run it:  npx playwright test --reporter=list <scratch-file>", WHITE),
         ("", WHITE),
         ("If it uses a CSS selector, ask: what ARIA role does this element have?", MUTED)])

    # ── 26. Debug Workflow ────────────────────────────────────────────────────
    sl = s()
    slide_concept(sl, "CH 2", BLUE, "Debug Workflow — Traces + Inspector",
        [("When a test fails:", MUTED),
         ("", WHITE),
         ("Step 1 — Run with trace:", WHITE),
         ("  npx playwright test --headed --trace=on --reporter=list", BGREEN),
         ("", WHITE),
         ("Step 2 — Open the trace:", WHITE),
         ("  npx playwright show-trace test-results/.../trace.zip", BGREEN),
         ("  Each action has: DOM snapshot before/after, network, console", MUTED),
         ("", WHITE),
         ("Step 3 — Pause at failure with Inspector:", WHITE),
         ("  PWDEBUG=1 npx playwright test", BGREEN),
         ("  Inspector opens at page.pause() — step through DOM live", MUTED),
         ("", WHITE),
         ("Step 4 — Use /pw-debug skill:", WHITE),
         ("  /pw-debug tests/auth.spec.ts", PURPLE),
         ("  Claude runs headed + trace, diagnoses root cause, proposes fix", MUTED)])

    # ── 27. Ch2 Discussion ────────────────────────────────────────────────────
    sl = s()
    slide_discussion(sl, "CH 2", BLUE, [
        "ActionWrapper.snapshot() returns both the accessibility tree and a screenshot. "
        "When would an LLM need the screenshot, and when is the a11y tree alone sufficient?",
        "navigateWithRetry multiplies delay by attempt number (retryDelayMs * attempt). "
        "Is this true exponential back-off? What would genuine exponential look like, "
        "and when would you prefer it?",
        "ActionResult captures durationMs for every action. Name two concrete ways "
        "an agent could use this timing data to improve future runs."])

    # ── 28. Ch2 Workshop Tasks ────────────────────────────────────────────────
    sl = s()
    slide_tasks(sl, "CH 2", BLUE, [
        ("Instrument action logs",
         "Add a console.log inside ActionWrapper.click() that prints: [click] selector durationMs ok:bool.\n"
         "Run a login flow. How many ms does each action take?"),
        ("Fix the back-off",
         "Change navigateWithRetry to use exponential back-off: retryDelayMs * 2^(attempt-1).\n"
         "Add a maxDelayMs cap of 10000ms."),
        ("codegen a login flow",
         "npx playwright codegen http://localhost:5173\n"
         "Perform login → check dashboard → products. Save. Run it. Zero CSS selectors?"),
        ("Debug a deliberate failure",
         "Change getByRole('button', {name:'Sign In'}) → locator('.nonexistent').\n"
         "Run with --trace=on. Open trace. Find the failure frame. Fix it."),
    ])

    # ── 29. LUNCH ─────────────────────────────────────────────────────────────
    sl = s()
    slide_break(sl, "🍽️  Lunch", "12:30 – 13:15  ·  45 minutes", GREEN)

    # ── 30. Half-time check-in ────────────────────────────────────────────────
    sl = s()
    slide_full_text(sl,
        "So far: you have the foundation.",
        "Agentic loop · ARIA selectors · Fault-tolerant sessions · CLI tooling\n\n"
        "Afternoon: MCP · Self-healing · Agents · Skills · CI/CD", BLUE)

    # ── 31. Ch3 Opener ────────────────────────────────────────────────────────
    sl = s()
    slide_chapter_opener(sl, "CH 3", "13:15",
        "Build Your Own MCP Server",
        "60 min  ·  Wire a browser as structured tools any LLM client can call", GREEN)

    # ── 32. MCP Architecture ─────────────────────────────────────────────────
    sl = s()
    bg(sl); topbar(sl, GREEN); chapter_label_top(sl, "CH 3", GREEN)
    tb(sl, 365760, 140000, W-730000, 440000, "MCP Architecture", 24, WHITE, bold=True)
    boxes = [
        (0, "Host", "Claude Code\nor VS Code\nCopilot", BLUE),
        (1, "Client", "MCP Client\n(built-in to\nthe host)", GREEN),
        (2, "Server", "Your\nserver.ts\n(stdio)", ORANGE),
        (3, "Browser", "Playwright\nChromium", PURPLE),
    ]
    bw = (W - 548640 - 3*120000) // 4
    bt = 680000; bh = 2800000
    for i, (idx, title, body, col) in enumerate(boxes):
        bx = 274320 + i*(bw+120000)
        rect(sl, bx, bt, bw, bh, CARD)
        rect(sl, bx, bt, bw, 60000, col)
        tb(sl, bx+80000, bt+120000, bw-160000, 300000, title, 16, col, bold=True, align=PP_ALIGN.CENTER)
        tb(sl, bx+80000, bt+480000, bw-160000, bh-600000, body, 13, MUTED, align=PP_ALIGN.CENTER)
        if i < 3:
            tb(sl, bx+bw+20000, bt+bh//2-100000, 80000, 200000, "→", 20, col, bold=True)
    multiline(sl, 274320, bt+bh+150000, W-548640, 1600000,
        [("Tool call flow:  Host prompts LLM  →  LLM emits tool_use JSON  →  Client routes to Server  →", WHITE),
         ("Server calls Playwright  →  result flows back  →  LLM sees result and plans next action.", WHITE),
         ("", WHITE),
         ("Config in .claude/settings.json:", MUTED),
         ('{ "mcpServers": { "web-detective": { "command": "npx", "args": ["ts-node", "examples/ch3-mcp/server.ts"] } } }', BGREEN)],
        13, font="Courier New" if False else None)

    # ── 33. Tool Schema Design ────────────────────────────────────────────────
    sl = s()
    slide_code(sl, "CH 3", GREEN, "Tool Schema Design — What the LLM Sees",
               "ch3-mcp/server.ts",
        [("server.tool('browser_click', {", WHITE),
         ("  description: 'Click an element. Prefer ARIA selectors.',", MUTED),
         ("  inputSchema: z.object({", WHITE),
         ("    selector:    z.string().describe('Playwright selector'),", WHITE),
         ("    description: z.string().describe('What you are clicking'),", WHITE),
         ("  }),", WHITE),
         ("}, async ({ selector, description }) => {", WHITE),
         ("  await page.locator(selector).click({ timeout: 5000 })", WHITE),
         ("  return { content: [{ type: 'text',", WHITE),
         ("    text: `Clicked: ${description}` }] }", WHITE),
         ("})", WHITE),
         ("", WHITE),
         ("// Rule: descriptions drive selector choice.", MUTED),
         ("// The LLM reads 'description' to decide WHICH selector to use.", MUTED),
         ("// Without a description, it will hallucinate a CSS path.", MUTED)])

    # ── 34. @playwright/mcp ───────────────────────────────────────────────────
    sl = s()
    slide_two_panel(sl, "CH 3", GREEN, "@playwright/mcp — Zero-Config Option",
        "Start in one command",
        [("npx @playwright/mcp@latest", BGREEN),
         ("", WHITE),
         ("Add to .claude/settings.json:", WHITE),
         ('{ "mcpServers": {', MUTED),
         ('    "playwright": {', MUTED),
         ('      "command": "npx",', MUTED),
         ('      "args": ["@playwright/mcp@latest"]', MUTED),
         ('    } } }', MUTED),
         ("", WHITE),
         ("For vision mode:", WHITE),
         ('"args": ["@playwright/mcp@latest", "--vision"]', MUTED)],
        "Built-in tools",
        [("browser_navigate(url)", BGREEN),
         ("browser_click(element)", BGREEN),
         ("browser_fill(element, value)", BGREEN),
         ("browser_snapshot()   ← a11y tree", BGREEN),
         ("browser_screenshot() ← PNG", BGREEN),
         ("browser_get_console_logs()", BGREEN),
         ("browser_tab_new / browser_tab_select", BGREEN),
         ("browser_network_requests()", BGREEN),
         ("", WHITE),
         ("No server.ts to write — plug in\nand start prompting.", MUTED)])

    # ── 35. Vision vs Snapshot ────────────────────────────────────────────────
    sl = s()
    slide_two_panel(sl, "CH 3", GREEN, "Snapshot Mode vs Vision Mode",
        "Snapshot (default) — prefer",
        [("browser_snapshot()  →  accessibility tree text", BGREEN),
         ("", WHITE),
         ("✓  No image tokens — much cheaper", WHITE),
         ("✓  Works headless in CI", WHITE),
         ("✓  Structured: roles, names, states", WHITE),
         ("✓  Forms, tables, lists: excellent", WHITE),
         ("✗  Can't see canvas or SVG charts", WHITE),
         ("", WHITE),
         ("Use for: login, products table,\nnavigation, form assertions", MUTED)],
        "Vision (--vision flag)",
        [("browser_screenshot()  →  PNG bytes", ORANGE),
         ("", WHITE),
         ("✓  Works on any visual element", WHITE),
         ("✓  Recharts / canvas / custom UI", WHITE),
         ("✗  Slow — screenshot + encode", WHITE),
         ("✗  Expensive image tokens", WHITE),
         ("✗  Sensitive to viewport / fonts", WHITE),
         ("", WHITE),
         ("Use for: dashboard charts,\nvisual regression, pixel checks", MUTED)])

    # ── 36. Ch3 Discussion ────────────────────────────────────────────────────
    sl = s()
    slide_discussion(sl, "CH 3", GREEN, [
        "The custom MCP server keeps a single shared AgentBrowserSession across all tool calls. "
        "What breaks if two AI clients connect simultaneously, and how would you fix it?",
        "browser_snapshot returns an accessibility tree while browser_screenshot returns a PNG. "
        "For a products-table assertion (count rows, read cell values), "
        "which tool is faster and cheaper — and why?",
        "A tool schema marks selector as required for browser_click. What error would the MCP "
        "server return if the LLM omits it, and where in the server code does that validation happen?"])

    # ── 37. Ch3 Workshop Tasks ────────────────────────────────────────────────
    sl = s()
    slide_tasks(sl, "CH 3", GREEN, [
        ("Wire @playwright/mcp to Claude Code",
         "Add mcpServers block to .claude/settings.json.\n"
         "Prompt: 'Log in with admin@shop.com / password123 and read the Total Revenue stat.'"),
        ("Compare snapshot vs vision",
         "Run the same prompt twice — once default, once with --vision.\n"
         "Which uses more tokens? Which answers the dashboard chart question better?"),
        ("Catch a console error via MCP",
         "Add console.error('chart render failed') to Dashboard.tsx useEffect.\n"
         "Prompt: 'Navigate to dashboard and check for console errors.'"),
        ("Debug a failed tool call",
         "Remove the selector field from browser_click's Zod schema.\n"
         "Observe the MCP error. Restore it. What was the validation message?"),
    ])

    # ── 38. Ch4 Opener ────────────────────────────────────────────────────────
    sl = s()
    slide_chapter_opener(sl, "CH 4", "14:15",
        "Self-Healing Selectors",
        "45 min  ·  Detect broken locators at runtime and auto-repair them", ORANGE)

    # ── 39. The Selector Rot Problem ─────────────────────────────────────────
    sl = s()
    slide_concept(sl, "CH 4", ORANGE, "The Selector Rot Problem",
        [("Every UI change potentially breaks tests.", MUTED),
         ("", WHITE),
         ("Designer renames .login-card → .auth-card:", WHITE),
         ("  20 test files need updating  —  discovered only when CI fails", RED),
         ("", WHITE),
         ("Developer reorders form fields:", WHITE),
         ("  nth-child(2) now points to the wrong input", RED),
         ("", WHITE),
         ("Team ships at velocity > maintenance bandwidth:", WHITE),
         ("  Test debt accumulates  →  suite becomes unreliable  →  ignored", RED),
         ("", WHITE),
         ("The self-healing approach:", MUTED),
         ("→ Keep a persistent LocatorStore — one selector per logical element", WHITE),
         ("→ On failure: snapshot the DOM, ask Claude for alternatives, validate each", WHITE),
         ("→ Write the winner back — next run uses the healed selector", BGREEN)])

    # ── 40. LocatorStore ─────────────────────────────────────────────────────
    sl = s()
    slide_code(sl, "CH 4", ORANGE, "LocatorStore — Persistent Locator Memory",
               "ch4-self-healing/locator-store.ts",
        [("interface LocatorEntry {", WHITE),
         ("  key:          string    // 'login.emailInput'", WHITE),
         ("  selector:     string    // current best", WHITE),
         ("  fallbacks:    string[]  // previous selectors, newest first", WHITE),
         ("  healCount:    number    // how many times auto-healed", WHITE),
         ("  lastVerified: string    // ISO timestamp of last passing run", WHITE),
         ("  lastHealed?:  string    // ISO timestamp of last heal event", WHITE),
         ("}", WHITE),
         ("", WHITE),
         ("// Seed data for the web-detective app:", MUTED),
         ("['login.emailInput',     'label:has-text(\"Email address\") + input'],", WHITE),
         ("['login.submitButton',   'button:has-text(\"Sign In\")'],", WHITE),
         ("['dashboard.statCards',  '.stat-card'],", WHITE),
         ("['products.searchInput', 'input[placeholder*=\"Search\"]'],", WHITE),
         ("", WHITE),
         ("store.heal('login.submitButton', 'button[type=\"submit\"]')", BGREEN),
         ("// → promotes new selector, demotes old to fallbacks, healCount += 1", MUTED)])

    # ── 41. The Healing Loop ──────────────────────────────────────────────────
    sl = s()
    bg(sl); topbar(sl, ORANGE); chapter_label_top(sl, "CH 4", ORANGE)
    tb(sl, 365760, 140000, W-730000, 440000,
       "The Self-Healing Loop — 6 Steps", 24, WHITE, bold=True)
    steps = [
        ("1", "Miss",    "Selector not found on page",           ORANGE),
        ("2", "Snapshot","page.locator('body').ariaSnapshot()",  BLUE),
        ("3", "Ask LLM", "Propose 3-5 alternatives ranked by stability", PURPLE),
        ("4", "Validate","Try each candidate on live DOM",       YELLOW),
        ("5", "Persist", "store.heal(key, winner)",              BGREEN),
        ("6", "Continue","Return winner — test proceeds",        BGREEN),
    ]
    step_w = (W - 548640 - 5*80000) // 6
    st = 680000; sh = 2600000
    for i, (num, title, desc, col) in enumerate(steps):
        sx = 274320 + i*(step_w+80000)
        rect(sl, sx, st, step_w, sh, CARD)
        rect(sl, sx, st, step_w, 50000, col)
        tb(sl, sx+40000, st+100000, step_w-80000, 350000,
           num, 26, col, bold=True, align=PP_ALIGN.CENTER)
        tb(sl, sx+40000, st+480000, step_w-80000, 350000,
           title, 14, WHITE, bold=True, align=PP_ALIGN.CENTER)
        tb(sl, sx+40000, st+860000, step_w-80000, sh-960000,
           desc, 11, MUTED, align=PP_ALIGN.CENTER)
        if i < 5:
            tb(sl, sx+step_w+10000, st+sh//2-100000, 60000, 200000,
               "→", 16, col, bold=True)
    # LLM call below steps
    rect(sl, 274320, st+sh+150000, W-548640, 2400000, CARD)
    tb(sl, 502920, st+sh+250000, W-1005840, 320000,
       "LLM Prompt (self-healer.ts)", 15, ORANGE, bold=True)
    multiline(sl, 502920, st+sh+620000, W-1005840, 1700000,
        [('system: "Given a broken selector and the current accessibility tree,", ', MUTED),
         ('"return a JSON array of 3-5 Playwright selectors. Prefer ARIA."', MUTED),
         ("", WHITE),
         ('user:  key: "login.submitButton"', WHITE),
         ('       brokenSelector: ".btn-login"', WHITE),
         ('       ariaTree: <page snapshot capped at 4000 chars>', WHITE),
         ("", WHITE),
         ('response: ["button:has-text(\\"Sign In\\")", "button[type=\\"submit\\"]", "getByRole(\'button\')"]', BGREEN)],
        12, font="Courier New")

    # ── 42. SelfHealingAgent Code ─────────────────────────────────────────────
    sl = s()
    slide_code(sl, "CH 4", ORANGE, "SelfHealingAgent — heal() Method",
               "ch4-self-healing/self-healer.ts",
        [("async heal(page, key, brokenSelector): Promise<HealResult> {", WHITE),
         ("  // Step 2: snapshot (cap at 4000 chars for token budget)", WHITE),
         ("  const domText = (await page.locator('body')", WHITE),
         ("    .ariaSnapshot()).slice(0, 4000)", MUTED),
         ("", WHITE),
         ("  // Step 3: ask LLM for alternatives", WHITE),
         ("  const candidates = await this.askLLMForCandidates(", WHITE),
         ("    key, brokenSelector, domText)", WHITE),
         ("", WHITE),
         ("  // Steps 4 + 5: validate each candidate", WHITE),
         ("  for (const candidate of candidates) {", WHITE),
         ("    const found = await page.locator(candidate)", WHITE),
         ("      .isVisible().catch(() => false)", WHITE),
         ("    if (found) {", WHITE),
         ("      this.store.heal(key, candidate)  // Step 5: persist", BGREEN),
         ("      return { healed: true, newSelector: candidate, candidates }", BGREEN),
         ("    }", WHITE),
         ("  }", WHITE),
         ("  return { healed: false, candidates }", ORANGE),
         ("}", WHITE)])

    # ── 43. Selector Stability ────────────────────────────────────────────────
    sl = s()
    slide_concept(sl, "CH 4", ORANGE, "Selector Stability Spectrum",
        [("Most stable  →  least stable:", MUTED),
         ("", WHITE),
         ("① ARIA role + accessible name        getByRole('button', {name:'Sign In'})", BGREEN),
         ("   Tied to user-visible semantics — survives CSS and DOM restructuring", MUTED),
         ("", WHITE),
         ("② Label text                          getByLabel('Email address')", BGREEN),
         ("   Tied to form labelling — survives style changes", MUTED),
         ("", WHITE),
         ("③ Visible text / placeholder          getByText('Sign In') / getByPlaceholder()", YELLOW),
         ("   Breaks on copy changes", MUTED),
         ("", WHITE),
         ("④ Semantic CSS class                  locator('.nav-link')", YELLOW),
         ("   Survives DOM restructuring, breaks on class renames", MUTED),
         ("", WHITE),
         ("⑤ data-testid attribute               locator('[data-testid=submit]')", YELLOW),
         ("   Stable but requires dev discipline to maintain", MUTED),
         ("", WHITE),
         ("⑥ Structural CSS / nth-child           .form > div:nth-child(2) > input", RED),
         ("   Breaks on any DOM reorder — never use in healer output", RED)])

    # ── 44. Ch4 Discussion ────────────────────────────────────────────────────
    sl = s()
    slide_discussion(sl, "CH 4", ORANGE, [
        "LocatorStore.heal() promotes the new selector to primary and demotes the old one to "
        "fallbacks. When should the healer try the fallbacks before asking the LLM for new candidates?",
        "The LLM prompt asks for '3-5 alternative selectors ordered by stability'. What makes a "
        "selector stable and how could you score stability without running the tests?",
        "healCount tracks how many times a locator has been replaced. At what threshold would you "
        "flag a locator for manual review instead of auto-healing it again?"])

    # ── 45. Ch4 Workshop Tasks ────────────────────────────────────────────────
    sl = s()
    slide_tasks(sl, "CH 4", ORANGE, [
        ("Seed and inspect the store",
         "const store = new LocatorStore('./locator-memory.json')\n"
         "for (const [k,v] of WEB_DETECTIVE_LOCATORS) store.register(k,v)\n"
         "store.printReport()  →  confirm all healCount === 0"),
        ("Break a selector and watch it heal",
         "In src/pages/Login.tsx rename .login-card → .auth-card  (keep both classes).\n"
         "Run healer.findElement(page, 'login.emailInput').  Watch heal() trigger."),
        ("Simulate heal() and trace state",
         "store.heal('login.submitButton', 'button[type=\"submit\"]')\n"
         "store.printReport()  →  old value in fallbacks, healCount 1, lastHealed set"),
        ("Add a toCsv() method",
         "Add toCsv(): string to LocatorStore — one line per entry: key,selector,healCount.\n"
         "Run npx ts-node -e '...' to print it. Commit the CSV for diff-friendly history."),
    ])

    # ── 46. Ch4.1 Opener ──────────────────────────────────────────────────────
    sl = s()
    slide_chapter_opener(sl, "CH 4.1", "15:00",
        "Playwright Agents: Planner, Generator, Healer",
        "30 min  ·  AI-generated test plans, code generation, and test-level healing",
        HORANGE)

    # ── 47. Three Built-In Agents ─────────────────────────────────────────────
    sl = s()
    bg(sl); topbar(sl, HORANGE); chapter_label_top(sl, "CH 4.1", HORANGE)
    tb(sl, 365760, 140000, W-730000, 440000,
       "Three Built-In Playwright Agents", 24, WHITE, bold=True)
    agents = [
        ("Planner",   HORANGE,
         "Explores the running app via the accessibility tree.\n"
         "Produces a structured Markdown test plan.\n"
         "Writes to specs/<name>.md",
         "Input: baseURL + routes\nOutput: human-readable plan"),
        ("Generator", BLUE,
         "Reads the Markdown plan.\n"
         "Emits runnable TypeScript Playwright tests.\n"
         "Verifies selectors live as it generates.",
         "Input: specs/<name>.md\nOutput: tests/<name>.spec.ts"),
        ("Healer",    BGREEN,
         "Runs the test suite.\n"
         "Replays failing steps on the live page.\n"
         "Rewrites broken test() blocks until green.",
         "Input: failing test file\nOutput: patched source file"),
    ]
    aw = (W - 548640 - 2*90000) // 3
    at = 680000; ah = H - at - 150000
    for i, (name, col, desc, io) in enumerate(agents):
        ax = 274320 + i*(aw+90000)
        rect(sl, ax, at, aw, ah, CARD)
        rect(sl, ax, at, aw, 60000, col)
        tb(sl, ax+120000, at+120000, aw-240000, 350000, name, 20, col, bold=True)
        tb(sl, ax+120000, at+520000, aw-240000, ah-1200000, desc, 13, WHITE)
        rect(sl, ax+60000, at+ah-700000, aw-120000, 560000, CARD2)
        tb(sl, ax+120000, at+ah-640000, aw-240000, 480000, io, 11, MUTED)

    # ── 48. Init Command + Pipeline ───────────────────────────────────────────
    sl = s()
    slide_two_panel(sl, "CH 4.1", HORANGE,
        "Initialise + Pipeline Flow",
        "One command to init",
        [("npx playwright init-agents --loop=claude", BGREEN),
         ("", WHITE),
         ("Generates agent definitions in .github/", WHITE),
         ("Works with:", WHITE),
         ("  • Claude Code  (--loop=claude)", WHITE),
         ("  • VS Code Copilot  (--loop=vscode)", WHITE),
         ("  • opencode  (--loop=opencode)", WHITE),
         ("", WHITE),
         ("Or run programmatically:", MUTED),
         ("npx ts-node examples/ch4.1-playwright-agents/planner.ts", BGREEN),
         ("→ Creates specs/ and tests/generated/", MUTED)],
        "Pipeline",
        [("Planner explores app", HORANGE),
         ("    ↓", MUTED),
         ("specs/web-detective.md", WHITE),
         ("    ↓", MUTED),
         ("Generator emits test code", BLUE),
         ("    ↓", MUTED),
         ("tests/generated/web-detective.spec.ts", WHITE),
         ("    ↓", MUTED),
         ("npx playwright test", MUTED),
         ("    ↓  (failures?)", MUTED),
         ("Healer rewrites test() blocks", BGREEN),
         ("    ↓", MUTED),
         ("Re-run  →  green", BGREEN)])

    # ── 49. Ch4 vs Ch4.1 ─────────────────────────────────────────────────────
    sl = s()
    slide_two_panel(sl, "CH 4.1", HORANGE, "Ch4 vs Ch4.1 — Healing at Two Levels",
        "Ch4  ·  Selector-Level Healing",
        [("Target: a broken selector string", ORANGE),
         ("Source: LocatorStore (locator-memory.json)", WHITE),
         ("Trigger: element.isVisible() returns false", WHITE),
         ("Output: new selector string written to store", WHITE),
         ("", WHITE),
         ("Best for:", MUTED),
         ("→ Ongoing test maintenance", WHITE),
         ("→ Production test suites", WHITE),
         ("→ Gradual UI drift over time", WHITE),
         ("", WHITE),
         ("Does NOT rewrite test source files.", MUTED)],
        "Ch4.1  ·  Test-Block Healing",
        [("Target: a failing test() block", HORANGE),
         ("Source: generated test source files", WHITE),
         ("Trigger: test exits with status: failed", WHITE),
         ("Output: patched .spec.ts written to disk", WHITE),
         ("", WHITE),
         ("Best for:", MUTED),
         ("→ First-run of AI-generated tests", WHITE),
         ("→ Bulk test generation + cleanup", WHITE),
         ("→ Large-scale migrations", WHITE),
         ("", WHITE),
         ("Does NOT update LocatorStore.", MUTED)])

    # ── 50. Ch4.1 Discussion ──────────────────────────────────────────────────
    sl = s()
    slide_discussion(sl, "CH 4.1", HORANGE, [
        "The Planner runs a seed test before navigating the app. "
        "What does the seed test actually bootstrap — and what would happen if the "
        "Planner navigated to http://localhost:5173 without it?",
        "The Generator is explicitly allowed to emit tests that contain errors, "
        "leaving them for the Healer to fix. Why is imperfect generation preferable to "
        "requiring the Generator to produce passing tests on the first attempt?",
        "A button is renamed from 'Submit' to 'Save' after the tests were generated. "
        "Ch4 SelfHealingAgent and Ch4.1 TestHealerAgent are both available. "
        "Which is the right tool, and why can't the other one help?"])

    # ── 51. Ch4.1 Workshop Tasks ──────────────────────────────────────────────
    sl = s()
    slide_tasks(sl, "CH 4.1", HORANGE, [
        ("Init and inspect agents",
         "npx playwright init-agents --loop=claude\n"
         "Open .github/ and read the generated agent definitions."),
        ("Run the Planner",
         "npx ts-node examples/ch4.1-playwright-agents/planner.ts\n"
         "Review specs/web-detective.md — does it cover login, dashboard, and products?"),
        ("Break a test and heal it",
         "In tests/generated/web-detective.spec.ts change one getByRole to locator('.gone').\n"
         "Run TestHealerAgent. Observe the multi-round repair loop."),
        ("Use /pw-plan skill",
         "/pw-plan web-detective\n"
         "Compare the skill-generated tests with the programmatic ones."),
    ])

    # ── 52. BREAK ─────────────────────────────────────────────────────────────
    sl = s()
    slide_break(sl, "☕  Break", "15:30 – 15:45  ·  15 minutes", HORANGE)

    # ── 53. Ch5 Opener ────────────────────────────────────────────────────────
    sl = s()
    slide_chapter_opener(sl, "CH 5", "15:45",
        "Building Your Own Custom AI Agent",
        "45 min  ·  Observe → Plan → Act → Verify → Learn with Ollama + DeepSeek-R1", BGREEN)

    # ── 54. WebTestAgent Architecture ────────────────────────────────────────
    sl = s()
    slide_concept(sl, "CH 5", BGREEN, "WebTestAgent — Architecture Overview",
        [("Three files. Three responsibilities.", MUTED),
         ("", WHITE),
         ("agent.ts     The reasoning loop: calls the local model, executes tools, collects results", WHITE),
         ("tools.ts     The tool registry: the JSON actions the model emits to call Playwright", WHITE),
         ("prompts.ts   The system + task prompts: defines identity and goal per run", WHITE),
         ("", WHITE),
         ("The loop (agent.ts):", MUTED),
         ("  1. Send goal + conversation history to DeepSeek-R1 via Ollama", WHITE),
         ("  2. Model returns ONE JSON action { tool, input }", WHITE),
         ("  3. Execute that tool against the live browser", WHITE),
         ("  4. Append the tool result to conversation history", WHITE),
         ("  5. Repeat until done() is called or 15 turns reached", WHITE),
         ("", WHITE),
         ("Key design choices:", MUTED),
         ("  • done() is a tool — not a return — so the model controls termination", WHITE),
         ("  • 15-turn limit prevents infinite loops if goal is impossible", WHITE),
         ("  • No native tool-calling → a prompt-based JSON action protocol drives dispatch", WHITE)])

    # ── 55. The Reasoning Loop Code ───────────────────────────────────────────
    sl = s()
    slide_code(sl, "CH 5", BGREEN, "The 15-Turn Reasoning Loop",
               "ch5-custom-agent/agent.ts",
        [("for (let turn = 0; turn < 15; turn++) {", WHITE),
         ("  const reply = await chat({", WHITE),
         ("    model:    process.env.WORKSHOP_MODEL ?? 'deepseek-r1:8b',", WHITE),
         ("    system:   SYSTEM_PROMPT,", WHITE),
         ("    messages,", WHITE),
         ("  })  // shared/ollama.ts strips <think>…</think> for us", BGREEN),
         ("  messages.push({ role: 'assistant', content: reply })", WHITE),
         ("", WHITE),
         ("  // Model returns ONE JSON action per turn: { tool, input }", MUTED),
         ("  const { tool, input } = JSON.parse(reply)", WHITE),
         ("  if (tool === 'done') return parseResult(input)", BGREEN),
         ("  const result = await this.execute(tool, input)", WHITE),
         ("", WHITE),
         ("  messages.push({ role: 'user',", WHITE),
         ("    content: JSON.stringify({ tool, result }) })", WHITE),
         ("}", WHITE)])

    # ── 56. Tool Registry ─────────────────────────────────────────────────────
    sl = s()
    slide_code(sl, "CH 5", BGREEN, "Tool Registry — What the Agent Can Do",
               "ch5-custom-agent/tools.ts",
        [("export type ToolName =", WHITE),
         ("  | 'navigate'        // go to a URL or path", WHITE),
         ("  | 'click'           // click element (selector + description required)", WHITE),
         ("  | 'fill'            // type into an input", WHITE),
         ("  | 'snapshot'        // ariaSnapshot — CALL FIRST before acting", BGREEN),
         ("  | 'assert_url'      // URL contains expected string", WHITE),
         ("  | 'assert_visible'  // element is visible on page", WHITE),
         ("  | 'assert_text'     // element contains expected text", WHITE),
         ("  | 'screenshot'      // take PNG when visual confirmation needed", WHITE),
         ("  | 'done'            // signal goal achieved/failed", BGREEN),
         ("", WHITE),
         ("// Every tool returns a string result consumed by the next turn.", MUTED),
         ("// 'done' is special: it short-circuits the loop immediately.", MUTED),
         ("", WHITE),
         ("// Adding a tool = add to ToolName, AGENT_TOOLS array, execute() switch.", MUTED)])

    # ── 57. System Prompt Design ──────────────────────────────────────────────
    sl = s()
    slide_concept(sl, "CH 5", BGREEN, "System Prompt Engineering — Rules That Matter",
        [("Always call snapshot() first.", BGREEN),
         ("  → Forces the agent to observe before acting. Prevents selector hallucination.", MUTED),
         ("", WHITE),
         ("Prefer ARIA selectors: getByRole, getByLabel, getByText, getByPlaceholder.", WHITE),
         ("  → Stable across style changes. Readable in logs.", MUTED),
         ("", WHITE),
         ("Use done() as soon as the goal is achieved or definitively failed.", WHITE),
         ("  → Without this rule the agent loops until the turn limit.", MUTED),
         ("", WHITE),
         ("Do not loop more than 15 tool calls.", WHITE),
         ("  → Hard guard — prevents runaway local inference / runtime.", MUTED),
         ("", WHITE),
         ("Assertions must produce binary outcomes: passed or failed, with evidence.", WHITE),
         ("  → Forces the agent to be explicit. Easy to parse in reporter.", MUTED),
         ("", WHITE),
         ("done() summary must include each assertion, the actual value observed,", WHITE),
         ("and whether it matched the expectation.", WHITE)])

    # ── 58. Prompt Caching ────────────────────────────────────────────────────
    sl = s()
    slide_concept(sl, "CH 5", BGREEN, "The JSON Action Protocol — No Native Tool Calls",
        [("DeepSeek-R1 has no native tool-calling API. We drive it with a prompt contract.", MUTED),
         ("Each turn the model must return exactly ONE JSON object — nothing else:", WHITE),
         ("  { \"tool\": \"click\", \"input\": { \"selector\": \"...\", \"description\": \"...\" } }", BGREEN),
         ("The loop parses it, dispatches to execute(tool, input), feeds the result back.", WHITE),
         ("", WHITE),
         ("Why it works:", MUTED),
         ("  • The system prompt lists every tool + its input schema as the contract.", WHITE),
         ("  • The model reasons in <think>…</think>; shared/ollama.ts strips it before parse.", WHITE),
         ("  • One action per turn keeps parsing deterministic and the loop auditable.", WHITE),
         ("", WHITE),
         ("Local warm-up, not billing:", MUTED),
         ("  • Runs on your machine via Ollama — no API key, no per-token cost.", BGREEN),
         ("  • Turn 1 is slow while the model loads into memory; later turns are faster.", WHITE),
         ("  • WORKSHOP_MODEL overrides the model (default deepseek-r1:8b).", WHITE)])

    # ── 59. Ch5 Discussion ────────────────────────────────────────────────────
    sl = s()
    slide_discussion(sl, "CH 5", BGREEN, [
        "The system prompt says 'Always call snapshot() first — never guess selectors.' "
        "What failure mode does this guard against, and what would you observe in the "
        "tool call log if the agent ignored this rule?",
        "done() is a regular tool in the registry rather than a hard loop exit. "
        "What is the advantage of this design, and what would break if you replaced "
        "it with a simple return from inside loop()?",
        "DeepSeek-R1 has no native tool-calling, so the agent uses a prompt-based JSON "
        "action protocol — one { tool, input } object per turn. What are the failure modes "
        "of parsing model output this way, and how would you make the loop robust to them?"])

    # ── 60. Ch5 Workshop Tasks ────────────────────────────────────────────────
    sl = s()
    slide_tasks(sl, "CH 5", BGREEN, [
        ("Run the agent",
         "npx ts-node examples/ch5-custom-agent/agent.ts\n"
         "Watch the per-step log. How many snapshot() calls in the full e-commerce run?"),
        ("Add a logout task prompt",
         "In prompts.ts add tasks.logout(). Goal: click Logout → assert URL has /login.\n"
         "Run agent.run(tasks.logout()). Does it pass?"),
        ("Add a get_text tool",
         "Add 'get_text' to ToolName and AGENT_TOOLS. Wire it to locator.innerText().\n"
         "Run fullEcommerce — does the agent use it to read stat card values?"),
        ("Measure warm-up time",
         "Log the wall-clock duration of each chat() call.\n"
         "Confirm turn 1 is slowest (model load) and later turns are faster."),
    ])

    # ── 61. Ch6 Opener ────────────────────────────────────────────────────────
    sl = s()
    slide_chapter_opener(sl, "CH 6", "16:30",
        "Creating Effective Testing Skills",
        "30 min  ·  Reusable, parameterised AI behaviours scoped to QA", PURPLE)

    # ── 62. What is a Skill? ──────────────────────────────────────────────────
    sl = s()
    slide_concept(sl, "CH 6", PURPLE, "Skills — Reusable Agent Behaviours",
        [("A skill is a SKILL.md file that gives an AI agent a specialised capability.", MUTED),
         ("Invoked with  /skill-name [arguments]  in Claude Code or VS Code.", WHITE),
         ("", WHITE),
         ("A well-crafted skill is the difference between:", WHITE),
         ("  an agent that works once  vs  one that works every time", WHITE),
         ("", WHITE),
         ("What makes a skill effective:", MUTED),
         ("  ✓ Clear argument-hint — users know what to pass", WHITE),
         ("  ✓ allowed-tools — pre-approves only the tools it needs", WHITE),
         ("  ✓ Specific output contract — 'list failing tests, then error, then one-line fix'", WHITE),
         ("  ✓ Guard rails — 'never modify test files unless user confirms'", WHITE),
         ("  ✓ Anchored to project paths — refers to actual fixtures, page objects", WHITE),
         ("", WHITE),
         ("Project skills in this repo:", MUTED),
         ("  /pw-run   /pw-debug   /pw-new-test   /pw-page-object   /pw-self-heal   /pw-plan", PURPLE)])

    # ── 63. SKILL.md Anatomy ──────────────────────────────────────────────────
    sl = s()
    slide_code(sl, "CH 6", PURPLE, "Anatomy of a SKILL.md",
               ".claude/skills/pw-run/SKILL.md",
        [("---", MUTED),
         ("name:          pw-run", WHITE),
         ("description:   Run Playwright tests. Pass file, grep pattern, or leave blank.", WHITE),
         ("argument-hint: \"[test-file-or-grep-pattern]\"   # shown in IDE autocomplete", PURPLE),
         ("allowed-tools: Bash(npx playwright *)           # pre-approved, no prompt", PURPLE),
         ("---", MUTED),
         ("", WHITE),
         ("Run the Playwright test suite. Arguments: $ARGUMENTS", WHITE),
         ("", WHITE),
         ("## Steps", WHITE),
         ("", WHITE),
         ("1. Determine scope from $ARGUMENTS:", WHITE),
         ("   - Empty  → npx playwright test --reporter=list", MUTED),
         ("   - File   → npx playwright test $ARGUMENTS --reporter=list", MUTED),
         ("   - String → npx playwright test --grep \"$ARGUMENTS\"", MUTED),
         ("", WHITE),
         ("2. Run command and capture output.", WHITE),
         ("3. Report: total passed/failed/skipped + duration.", WHITE),
         ("4. For failures: quote error + suggest fix by category.", WHITE),
         ("Never modify test files unless the user explicitly asks.", RED)])

    # ── 64. Project Skills Table ──────────────────────────────────────────────
    sl = s()
    slide_table(sl, "CH 6", PURPLE, "Project Skills — Complete Inventory",
        ["Skill", "Purpose", "allowed-tools"],
        [["/pw-run",         "Run full suite or filtered subset",           "Bash(npx playwright *)"],
         ["/pw-debug",       "Headed run + trace + root-cause diagnosis",   "Bash(npx playwright *)"],
         ["/pw-new-test",    "Scaffold spec using fixtures + page objects",  "Bash, Read, Write"],
         ["/pw-page-object", "Generate typed POM from React source",        "Read, Write"],
         ["/pw-self-heal",   "Detect broken selectors, propose ARIA fixes", "Bash, Read, Edit"],
         ["/pw-plan",        "Run Planner → Generator pipeline",            "Bash(npm *), Bash(npx *)"]])

    # ── 65. Prompt Engineering in Skills ─────────────────────────────────────
    sl = s()
    slide_two_panel(sl, "CH 6", PURPLE, "Prompt Engineering Inside Skills",
        "What works",
        [("Be specific about output format:", WHITE),
         ("  'list failing test titles, then the", WHITE),
         ("  error, then a one-line fix'", WHITE),
         ("  beats 'explain the failure'", MUTED),
         ("", WHITE),
         ("Anchor to project conventions:", WHITE),
         ("  Reference actual file paths,\n  fixture names, page object patterns", WHITE),
         ("", WHITE),
         ("Guard rails that matter:", WHITE),
         ("  'Never modify test files unless\n  the user confirms'", WHITE),
         ("  'Do not run npm install'", WHITE)],
        "What to avoid",
        [("Vague descriptions:", RED),
         ("  'help with tests'  →  LLM drifts", MUTED),
         ("", WHITE),
         ("Over-broad allowed-tools:", RED),
         ("  * (all tools) means the agent\n  can do anything", MUTED),
         ("", WHITE),
         ("Missing output contract:", RED),
         ("  Without a format spec the agent\n  improvises — inconsistent output", MUTED),
         ("", WHITE),
         ("No argument-hint:", RED),
         ("  Collaborators won't know what\n  to pass as $ARGUMENTS", MUTED)])

    # ── 66. Ch6 Discussion ────────────────────────────────────────────────────
    sl = s()
    slide_discussion(sl, "CH 6", PURPLE, [
        "The /pw-debug skill injects live test output via the shell before the LLM sees "
        "the prompt. Why does running the tests inside the skill (rather than pasting "
        "output manually) produce better diagnoses?",
        "allowed-tools in a skill's frontmatter restricts which tools the agent may call. "
        "What is the security and correctness argument for limiting /pw-new-test to only "
        "Read, Edit, Write, and Bash?",
        "A skill with no argument-hint field still works when invoked with an argument. "
        "What is the purpose of argument-hint, and who specifically benefits from it?"])

    # ── 67. Ch6 Workshop Tasks ────────────────────────────────────────────────
    sl = s()
    slide_tasks(sl, "CH 6", PURPLE, [
        ("Run an existing skill",
         "/pw-run tests/auth.spec.ts\n"
         "Observe: does the skill summarise pass/fail? Does it suggest a fix on failure?"),
        ("Write a debug-test skill from scratch",
         "Create .claude/skills/debug-test/SKILL.md with: argument-hint, allowed-tools,\n"
         "headed run step, root-cause categories, output contract. Invoke it."),
        ("Audit allowed-tools",
         "Check each of the 6 project skills. Is any allowed-tools broader than needed?\n"
         "Reduce scope on one skill and verify it still works correctly."),
        ("Add cross-reference to pw-self-heal",
         "Read .claude/skills/pw-self-heal/SKILL.md scope note.\n"
         "Does it correctly redirect generated test failures to ch4.1 TestHealerAgent?"),
    ])

    # ── 68. Ch7 Opener ────────────────────────────────────────────────────────
    sl = s()
    slide_chapter_opener(sl, "CH 7", "17:00",
        "Implementing an Agent & Running on CI",
        "60 min  ·  Structured output, annotations, parallel matrix, post-mortem reporting",
        BLUE)

    # ── 69. CI Contract ───────────────────────────────────────────────────────
    sl = s()
    slide_concept(sl, "CH 7", BLUE, "The CI Contract — Exit Codes + Critical Flag",
        [("Every agent run maps to one of three exit codes:", MUTED),
         ("", WHITE),
         ("  0   All critical scenarios passed  (pipeline continues)", BGREEN),
         ("  1   At least one critical scenario failed  (pipeline fails)", RED),
         ("  2   Configuration error (missing API key, bad scenario name)", ORANGE),
         ("", WHITE),
         ("The critical flag per scenario:", MUTED),
         ("  critical: true   →  failure = exit 1, ::error annotation in PR diff", RED),
         ("  critical: false  →  failure = exit 0, ::warning annotation only", YELLOW),
         ("", WHITE),
         ("Why separate critical from optional:", WHITE),
         ("  Some checks are exploratory (visual, performance, a11y) — they should warn", WHITE),
         ("  but never block a deploy. Mixing them inflates noise and causes bypass.", WHITE),
         ("", WHITE),
         ("Rule of thumb:", MUTED),
         ("  Start new scenarios as critical: false.", WHITE),
         ("  Promote to critical: true once they've run cleanly for 5+ consecutive builds.", WHITE)])

    # ── 70. CIReporter ────────────────────────────────────────────────────────
    sl = s()
    slide_code(sl, "CH 7", BLUE, "CIReporter — Four Output Channels",
               "ch7-agent-ci/reporter.ts",
        [("class CIReporter {", WHITE),
         ("  record(result: ScenarioResult): void {", WHITE),
         ("    // Channel 1: workflow annotation → inline in PR diff", WHITE),
         ("    this.annotation(result.critical ? 'error' : 'warning',", WHITE),
         ("      `Agent scenario failed: ${result.name}`, result.summary)", BLUE),
         ("  }", WHITE),
         ("", WHITE),
         ("  async writeStepSummary(): Promise<void> {", WHITE),
         ("    // Channel 2: Markdown table in $GITHUB_STEP_SUMMARY", WHITE),
         ("    fs.appendFileSync(process.env.GITHUB_STEP_SUMMARY, markdownTable)", BLUE),
         ("  }", WHITE),
         ("", WHITE),
         ("  setOutputs(): void {", WHITE),
         ("    // Channel 3: job outputs → consumed by downstream jobs", WHITE),
         ("    this.setOutput('passed', String(allCriticalPassed))", BLUE),
         ("  }", WHITE),
         ("", WHITE),
         ("  writeReport(path): void {", WHITE),
         ("    // Channel 4: agent-report.json → uploaded as artifact", WHITE),
         ("    fs.writeFileSync(path, JSON.stringify(report, null, 2))", BLUE),
         ("  }", WHITE),
         ("}", WHITE)])

    # ── 71. GitHub Actions Annotations ────────────────────────────────────────
    sl = s()
    slide_concept(sl, "CH 7", BLUE, "GitHub Actions Workflow Commands",
        [("Emitted to stdout — GitHub Actions intercepts and renders them:", MUTED),
         ("", WHITE),
         ("::error title=<title>::<message>", RED),
         ("  → Red annotation inline in the PR diff file view", MUTED),
         ("", WHITE),
         ("::warning title=<title>::<message>", YELLOW),
         ("  → Yellow annotation — visible but doesn't block merge", MUTED),
         ("", WHITE),
         ("::notice title=<title>::<message>", MUTED),
         ("  → Neutral — useful for informational annotations", MUTED),
         ("", WHITE),
         ("Outside CI (local run): CIReporter prints to stdout instead.", WHITE),
         ("  ✗ [Agent scenario failed: login-flow]", RED),
         ("  ⚠ [Agent scenario failed: logout-flow]", YELLOW),
         ("", WHITE),
         ("Message escaping required:", MUTED),
         ("  % → %25   \\r → %0D   \\n → %0A   : → %3A   , → %2C", WHITE)])

    # ── 72. Scenario Registry ────────────────────────────────────────────────
    sl = s()
    slide_code(sl, "CH 7", BLUE, "Scenario Registry — Decoupled from Pipeline YAML",
               "ch7-agent-ci/scenarios.ts",
        [("interface Scenario {", WHITE),
         ("  name:     string   // matrix key — matches GitHub Actions job ID", WHITE),
         ("  task:     string   // goal string passed to WebTestAgent.run()", WHITE),
         ("  critical: boolean  // false = warn-only, true = fails the build", WHITE),
         ("}", WHITE),
         ("", WHITE),
         ("export const SCENARIOS: Scenario[] = [", WHITE),
         ("  { name: 'auth-redirect',  task: tasks.authRedirect(),         critical: true  },", WHITE),
         ("  { name: 'login-flow',     task: tasks.loginFlow(),            critical: true  },", WHITE),
         ("  { name: 'product-search', task: tasks.productSearch('TV',2),  critical: true  },", WHITE),
         ("  { name: 'full-ecommerce', task: tasks.fullEcommerce(),        critical: true  },", WHITE),
         ("]", WHITE),
         ("", WHITE),
         ("// Adding a new scenario = add one object here.", MUTED),
         ("// No YAML changes needed — the matrix reads SCENARIOS at runtime.", MUTED)])

    # ── 73. Pipeline Architecture ─────────────────────────────────────────────
    sl = s()
    bg(sl); topbar(sl, BLUE); chapter_label_top(sl, "CH 7", BLUE)
    tb(sl, 365760, 140000, W-730000, 440000,
       "Pipeline Architecture", 24, WHITE, bold=True)
    # Main column
    jobs = [
        ("playwright (every push/PR)",       "~2 min fast gate, no API cost",              BGREEN),
        ("locator-health (every push/PR)",   "Seed store, probe selectors on live app",    ORANGE),
        ("agentic matrix (main only)",       "Parallel jobs — one per scenario",           BLUE),
    ]
    jw = W - 548640; jh = 1100000
    jt = 680000
    for i, (title, sub, col) in enumerate(jobs):
        jy = jt + i*(jh+80000)
        rect(sl, 274320, jy, jw, jh, CARD)
        rect(sl, 274320, jy, 8000, jh, col)
        tb(sl, 420000, jy+140000, jw-300000, 400000, title, 16, WHITE, bold=True)
        tb(sl, 420000, jy+580000, jw-300000, 400000, sub, 13, MUTED)
    # matrix inner
    matrix_top = jt + 2*(jh+80000)
    matrix_scenarios = ["auth-redirect", "login-flow", "product-search", "full-ecommerce"]
    sw = (jw - 200000 - 3*60000) // 4
    for i, sc in enumerate(matrix_scenarios):
        sx = 274320 + 100000 + i*(sw+60000)
        rect(sl, sx, matrix_top+160000, sw, jh-320000, CARD2)
        tb(sl, sx+60000, matrix_top+300000, sw-120000, jh-640000,
           sc, 11, BLUE, align=PP_ALIGN.CENTER)
    multiline(sl, 274320, matrix_top+jh+80000, jw, 800000,
        [("fail-fast: false  →  all matrix jobs finish even if one fails.", WHITE),
         ("Each job uploads: agent-report-<name>.json + trace.zip", MUTED),
         ("Cost gate: agentic matrix runs only on main — never on feature branches.", ORANGE)],
        13)

    # ── 74. Part A — CI Reporter ─────────────────────────────────────────────
    sl = s()
    slide_concept(sl, "CH 7  · Part A", BLUE, "Part A — CI Reporter Walkthrough",
        [("npx ts-node examples/ch7-agent-ci/agent-runner.ts --scenario login-flow", BGREEN),
         ("", WHITE),
         ("Without GITHUB_ACTIONS set, reporter prints to stdout:", WHITE),
         ("  ✓ [Agent scenario passed: login-flow]  Summary: ...", BGREEN),
         ("  ✗ [Agent scenario failed: product-search]  Error: ...", RED),
         ("", WHITE),
         ("After the run, open agent-report.json:", MUTED),
         ('  { "timestamp": "...", "passed": true, "totalDurationMs": 8420,', WHITE),
         ('    "scenarios": [{ "name": "login-flow", "passed": true,', WHITE),
         ('      "toolCalls": 7, "durationMs": 3210, "critical": true }] }', WHITE),
         ("", WHITE),
         ("Hands-on:", MUTED),
         ("1. Add a new annotation type (notice) to CIReporter.annotation()", WHITE),
         ("2. Verify it prints the ✓ prefix locally", WHITE),
         ("3. Check how $GITHUB_STEP_SUMMARY is written (appendFileSync)", WHITE)])

    # ── 75. Part B — Scenario Runner ─────────────────────────────────────────
    sl = s()
    slide_concept(sl, "CH 7  · Part B", BLUE, "Part B — Scenario Runner",
        [("CLI interface:", MUTED),
         ("  npx ts-node examples/ch7-agent-ci/agent-runner.ts --list", BGREEN),
         ("  npx ts-node examples/ch7-agent-ci/agent-runner.ts --scenario login-flow", BGREEN),
         ("  npx ts-node examples/ch7-agent-ci/agent-runner.ts --all", BGREEN),
         ("", WHITE),
         ("Exit codes:", MUTED),
         ("  0   all critical scenarios passed", BGREEN),
         ("  1   at least one critical scenario failed", RED),
         ("  2   unknown --scenario name", ORANGE),
         ("", WHITE),
         ("Hands-on:", MUTED),
         ("1. Run --scenario product-search. Observe per-step log lines.", WHITE),
         ("2. Add a logout-flow scenario with critical: false", WHITE),
         ("3. Make it fail deliberately (impossible count). Run --all.", WHITE),
         ("4. Confirm exit code is 0 (not critical) vs 1 (critical).", WHITE),
         ("   echo $?  after the command", BGREEN)])

    # ── 76. Ch7 Discussion ────────────────────────────────────────────────────
    sl = s()
    slide_discussion(sl, "CH 7", BLUE, [
        "The pipeline uses fail-fast: false on the agentic matrix. Under what "
        "circumstances would you change this to fail-fast: true, "
        "and what is the cost of doing so?",
        "CIReporter.annotation() emits ::error for critical failures and ::warning "
        "for non-critical ones. GitHub renders these inline in PR diffs. "
        "Why would ::notice be a poor choice for a broken selector?",
        "The agentic matrix runs only on main, while standard Playwright tests run "
        "on every push and PR. What tradeoff does this architecture make, and how "
        "would you structure a middle ground (e.g. run agents on PRs to main only)?"])

    # ── 77. Ch7 Workshop Tasks ────────────────────────────────────────────────
    sl = s()
    slide_tasks(sl, "CH 7", BLUE, [
        ("Run a single scenario",
         "npx ts-node examples/ch7-agent-ci/agent-runner.ts --scenario product-search\n"
         "Open agent-report.json. How many tool calls did the scenario use?"),
        ("Flip critical and observe exit code",
         "Set product-search critical: false. Make it assert an impossible count.\n"
         "Run --all. echo $?  →  should be 0. Flip to true. echo $?  →  should be 1."),
        ("Add a logout scenario",
         "Add logout task to prompts.ts. Add { name: 'logout-flow', critical: false }.\n"
         "Run it. Does the agent call Logout button correctly?"),
        ("Walk the pipeline YAML",
         "Open .github/workflows/agentic-tests.yml. Trace: trigger → type-check → matrix.\n"
         "Find where the agent connects to Ollama and why the agentic matrix is gated to main (runtime cost)."),
    ])

    # ── 78. Full Architecture ─────────────────────────────────────────────────
    sl = s()
    bg(sl); topbar(sl, BLUE)
    tb(sl, 365760, 100000, W-730000, 440000, "The Complete Architecture", 26, WHITE, bold=True)
    tb(sl, 365760, 580000, W-730000, 300000,
       "How all seven chapters connect:", 14, MUTED)
    layers = [
        ("LLM Layer",       "DeepSeek-R1  ·  Ollama  ·  Local inference  ·  JSON actions", PURPLE),
        ("Agent Layer",     "WebTestAgent (Ch5)  ·  MCP Server (Ch3)  ·  Playwright Agents (Ch4.1)", BLUE),
        ("Skills Layer",    "/pw-run  /pw-debug  /pw-plan  /pw-self-heal  /pw-new-test (Ch6)", PURPLE),
        ("Healing Layer",   "SelfHealingAgent (Ch4)  ·  LocatorStore  ·  TestHealerAgent (Ch4.1)", ORANGE),
        ("Execution Layer", "ActionWrapper (Ch2)  ·  BrowserSession  ·  Playwright CLI", BGREEN),
        ("CI Layer",        "CIReporter (Ch7)  ·  Scenario registry  ·  GitHub Actions matrix", BLUE),
        ("App Layer",       "React app  ·  /login  /dashboard  /products  ·  localhost:5173", GREEN),
    ]
    lh = (H - 1000000 - 150000) // len(layers)
    for i, (label, detail, col) in enumerate(layers):
        ly = 1000000 + i*(lh+18000)
        rect(sl, 274320, ly, W-548640, lh, CARD)
        rect(sl, 274320, ly, 60000, lh, col)
        tb(sl, 420000, ly+lh//2-120000, 2800000, 240000, label, 13, col, bold=True)
        tb(sl, 3400000, ly+lh//2-120000, W-3800000, 240000, detail, 13, WHITE)

    # ── 79. Key Takeaways ─────────────────────────────────────────────────────
    sl = s()
    slide_takeaways(sl, [
        ("Agentic ≠ Magic",
         "Agents are grounded in reliable browser automation — intelligence lives in orchestration."),
        ("ARIA Selectors First",
         "getByRole / getByLabel survive CSS and DOM changes. Structural selectors don't."),
        ("MCP is the Glue",
         "Any LLM client can drive a real browser once it speaks the MCP tool protocol."),
        ("Self-Healing is a Strategy",
         "LocatorStore + LLM repair beats brittle selectors at scale. Requires review discipline."),
        ("Skills = Reusable Behaviour",
         "Tight allowed-tools + clear output contract = a skill that works every time."),
        ("Scenarios are the Unit of CI",
         "Each scenario has a critical flag, a name, and lives in a registry — not YAML."),
        ("CI/CD is the Finish Line",
         "Agents that can't emit structured output and annotations are prototypes, not solutions."),
    ])

    # ── 80. Where the Field Is Heading ────────────────────────────────────────
    sl = s()
    slide_concept(sl, "Closing", BLUE, "Where the Field Is Heading",
        [("Multi-agent QA teams", BLUE),
         ("  Planner, Generator, Healer, and Reviewer agents working in parallel pipelines.", MUTED),
         ("", WHITE),
         ("Visual regression agents", PURPLE),
         ("  Vision-mode agents that catch pixel regressions without fragile screenshot diffs.", MUTED),
         ("", WHITE),
         ("Compliance agents", ORANGE),
         ("  Agents that verify WCAG, GDPR, and OWASP compliance continuously.", MUTED),
         ("", WHITE),
         ("Test-as-specification", BGREEN),
         ("  High-level PRDs → agent generates tests → tests drive implementation.", MUTED),
         ("", WHITE),
         ("Human-in-the-loop QA", WHITE),
         ("  Agents triage failures, humans approve fixes — 10× the coverage with same headcount.", MUTED)])

    # ── 81. Resources ─────────────────────────────────────────────────────────
    sl = s()
    slide_concept(sl, "Closing", BLUE, "Resources",
        [("Playwright", BLUE),
         ("  playwright.dev  ·  playwright.dev/docs/test-agents", WHITE),
         ("", WHITE),
         ("Ollama", PURPLE),
         ("  ollama.com  ·  npmjs.com/package/ollama", WHITE),
         ("", WHITE),
         ("DeepSeek-R1", HORANGE),
         ("  ollama.com/library/deepseek-r1", WHITE),
         ("", WHITE),
         ("MCP", GREEN),
         ("  modelcontextprotocol.io  ·  npmjs.com/package/@playwright/mcp", WHITE),
         ("", WHITE),
         ("GitHub Actions", BGREEN),
         ("  docs.github.com/en/actions/writing-workflows/workflow-syntax", WHITE),
         ("", WHITE),
         ("Claude Code", ORANGE),
         ("  docs.anthropic.com/en/docs/claude-code  ·  /help inside Claude Code", WHITE),
         ("", WHITE),
         ("This repo", MUTED),
         ("  github.com/[your-org]/web-detective  ·  npm install  ·  npm run dev", WHITE)])

    # ── 82. Q&A ───────────────────────────────────────────────────────────────
    sl = s()
    bg(sl)
    rect(sl, 0, 0, W, H//2, BLUE)
    rect(sl, 0, H//2, W, H - H//2, CARD)
    tb(sl, 0, H//2 - 600000, W, 500000,
       "Questions?", 54, BLACK, bold=True, align=PP_ALIGN.CENTER)
    tb(sl, 0, H//2 + 250000, W, 400000,
       "Web Detective Workshop", 22, WHITE, bold=True, align=PP_ALIGN.CENTER)
    multiline(sl, 0, H//2 + 800000, W, 600000,
        [("playwright.dev  ·  ollama.com  ·  modelcontextprotocol.io", MUTED),
         ("admin@shop.com  /  password123  ·  http://localhost:5173", MUTED)],
        13, align=PP_ALIGN.CENTER)

    print(f"Built {len(prs.slides)} slides.")

# ── Entry point ───────────────────────────────────────────────────────────────

if __name__ == "__main__":
    prs = Presentation()
    prs.slide_width  = Emu(W)
    prs.slide_height = Emu(H)
    build(prs)
    prs.save("web-detective-workshop.pptx")
    print("Saved → web-detective-workshop.pptx")
